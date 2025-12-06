import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
import sys

# -----------------------------------------------
# 상수 정의
# -----------------------------------------------
CSV_PATH = Path("problems.csv")
PNG_DIR = Path("output_pngs") # convert_with_matplotlib.py의 결과 폴더
PPT_DIR = Path("output_ppts") # 완성된 PPT를 저장할 폴더

# -----------------------------------------------
# PPT 생성 메인 함수
# -----------------------------------------------
def create_presentation(concept_name):
    """
    특정 개념(concept)에 대한 문제들로 PPT 프레젠테이션을 생성합니다.
    """
    
    # 1. CSV 데이터 읽기
    try:
        df = pd.read_csv(CSV_PATH)
    except FileNotFoundError:
        print(f"[Error] CSV 파일을 찾을 수 없습니다: {CSV_PATH}")
        print("스크립트와 같은 위치에 problems.csv 파일이 있는지 확인하세요.")
        return
    except Exception as e:
        print(f"[Error] CSV 파일 읽기 오류: {e}")
        return

    # 2. 원하는 'concept'으로 데이터 필터링
    problems_df = df[df['concept'] == concept_name].copy()
    problems_df = problems_df.dropna(subset=['id', 'problem_text', 'problem_latex']) # 필수 정보 없는 행 제거

    if problems_df.empty:
        print(f"[Info] '{concept_name}'에 해당하는 문제를 찾을 수 없습니다.")
        return

    print(f"'{concept_name}' 개념에 대해 총 {len(problems_df)}개의 문제로 PPT 생성을 시작합니다...")

    # 3. PPT 프레젠테이션 객체 생성 (기본 16:9 비율)
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # 4. 제목 슬라이드 추가
    title_slide_layout = prs.slide_layouts[0] # 제목 슬라이드 레이아웃
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1] # 부제목
    
    title.text = f"수학 수업 자료: {concept_name}"
    subtitle.text = f"{len(problems_df)}개의 문제 수록"

    # 5. 문제 및 정답 슬라이드 추가
    problem_layout = prs.slide_layouts[5] # '제목만' 레이아웃 (내용은 직접 추가)

    for idx, row in problems_df.iterrows():
        # 파일명으로 사용할 ID 정규화 (convert_with_matplotlib.py와 동일한 로직)
        file_id = row.get("id")
        if pd.isna(file_id):
            file_id = f"idx_{idx}"
        else:
            try:
                file_id = int(file_id)
            except ValueError:
                file_id = f"id_{file_id}"
        
        problem_text = str(row.get("problem_text", "")).strip()
        problem_png_path = PNG_DIR / f"problem_{file_id}.png"
        answer_png_path = PNG_DIR / f"answer_{file_id}.png"

        # --- 문제 슬라이드 생성 ---
        slide = prs.slides.add_slide(problem_layout)
        title_shape = slide.shapes.title
        title_shape.text = f"문제 - {problem_text}"
        
        # PNG 이미지 삽입 (가로 중앙 배치)
        if problem_png_path.exists():
            img_width_in = 12.0 # 16인치 너비에 12인치 크기로
            img_height_auto = True # 높이는 비율에 맞게 자동 조절
            pic_left = Inches((16.0 - img_width_in) / 2) # 중앙 정렬
            pic_top = Inches(1.5) # 제목 아래
            
            slide.shapes.add_picture(
                str(problem_png_path), 
                pic_left, pic_top, 
                width=Inches(img_width_in)
            )
        else:
            print(f"[Warning] 문제 이미지를 찾을 수 없음: {problem_png_path}")
            slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(1)).text = "문제 이미지 없음"

        # --- 정답 슬라이드 생성 (정답 PNG가 있을 경우에만) ---
        if answer_png_path.exists():
            slide = prs.slides.add_slide(problem_layout)
            title_shape = slide.shapes.title
            title_shape.text = f"정답 - {problem_text}"

            slide.shapes.add_picture(
                str(answer_png_path),
                pic_left, pic_top,
                width=Inches(img_width_in)
            )
        else:
            # 정답 이미지가 없으면 정답 슬라이드를 아예 만들지 않거나, 텍스트만 표시
            print(f"[Info] 정답 이미지가 없음: {answer_png_path}")

    # 6. PPT 파일 저장
    PPT_DIR.mkdir(exist_ok=True) # output_ppts 폴더 생성
    
    # 파일명 생성 (특수문자 제거)
    safe_filename = "".join(c for c in concept_name if c.isalnum() or c in (' ', '_')).rstrip()
    safe_filename = safe_filename.replace(' ', '_') + "_수업자료.pptx"
    
    output_path = PPT_DIR / safe_filename
    prs.save(output_path)
    
    print("\n🎉 PPT 생성 완료!")
    print(f"'{output_path}' 경로에 파일이 저장되었습니다.")

# -----------------------------------------------
# 스크립트 실행 (사용자 입력 처리)
# -----------------------------------------------
def get_user_choice():
    """CSV에서 가능한 개념 목록을 보여주고 사용자 선택을 받습니다."""
    try:
        df = pd.read_csv(CSV_PATH)
        # NaN 값을 제거하고, 고유한 concept 목록을 만듭니다.
        available_concepts = df['concept'].dropna().unique()
        
        if len(available_concepts) == 0:
            print("[Error] CSV 파일에서 'concept' 컬럼을 찾을 수 없거나 비어있습니다.")
            return None

        print("-" * 30)
        print("생성할 PPT의 단원(concept)을 선택하세요:")
        print("-" * 30)
        
        for i, concept in enumerate(available_concepts):
            print(f"  [{i + 1}] {concept}")
        
        print("-" * 30)
        
        while True:
            choice_str = input(f"번호 (1~{len(available_concepts)}) 또는 단원명을 입력하세요: ").strip()
            
            # 1. 번호로 입력했는지 확인
            try:
                choice_idx = int(choice_str)
                if 1 <= choice_idx <= len(available_concepts):
                    return available_concepts[choice_idx - 1]
                else:
                    print("잘못된 번호입니다. 다시 입력해주세요.")
            except ValueError:
                # 2. 이름으로 입력했는지 확인
                if choice_str in available_concepts:
                    return choice_str
                else:
                    print("정확한 단원명이 아닙니다. 다시 입력해주세요.")
                    
    except FileNotFoundError:
        print(f"[Error] CSV 파일을 찾을 수 없습니다: {CSV_PATH}")
        return None
    except Exception as e:
        print(f"[Error] CSV 파일 처리 중 오류 발생: {e}")
        return None

if __name__ == "__main__":
    # 0. PNG 파일이 준비되었는지 확인
    if not PNG_DIR.exists() or not any(PNG_DIR.iterdir()):
        print(f"[Warning] '{PNG_DIR}' 폴더가 비어있습니다.")
        print("먼저 'convert_with_matplotlib.py' 스크립트를 실행하여 PNG 이미지를 생성해야 합니다.")
        if input("무시하고 계속 진행하시겠습니까? (y/n): ").lower() != 'y':
            sys.exit("작업을 중단합니다.")

    # 1. 사용자에게 만들고 싶은 단원 선택받기
    selected_concept = get_user_choice()
    
    # 2. 사용자가 선택했을 경우에만 PPT 생성 실행
    if selected_concept:
        create_presentation(selected_concept)